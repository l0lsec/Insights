"""Extract plain text from uploaded documents so it can be fed to the LLM.

The compose "From File" tab lets users upload a document (PDF, PowerPoint,
Word, Excel, CSV, plain text, Markdown, or HTML) as source material for post
generation. Because the app is provider-agnostic (OpenAI / Anthropic / local
Ollama), we extract text server-side and reuse ``generate_posts_from_text``
rather than relying on any single provider's native file support.

Heavy parsers (``pypdf``, ``python-pptx``, ``python-docx``) are imported
lazily so a missing optional dependency only disables that one format instead
of breaking the whole app. Use :func:`supported_formats` to report which
extensions are actually usable in the current environment.
"""

from __future__ import annotations

import csv
import io
import os
import re
from typing import Callable, NamedTuple


class ExtractionError(Exception):
    """Raised when a document cannot be turned into usable text."""


# Cap extracted text so a huge upload can't exhaust memory or the prompt.
# ``generate_posts_from_text`` truncates further, but we keep a generous
# ceiling here and surface the real length to the caller.
MAX_EXTRACT_CHARS = 200_000

# Cap on how many bytes we'll read from an upload (mirrors Flask's
# MAX_CONTENT_LENGTH of 16MB; kept here so the extractor is usable standalone).
MAX_FILE_BYTES = 16 * 1024 * 1024


class FormatSpec(NamedTuple):
    label: str            # human-readable name for the UI
    extract: Callable      # fn(data: bytes) -> str
    requires: str | None   # importable module needed, or None if built-in only


def _truncate(text: str) -> str:
    text = text.strip()
    if len(text) > MAX_EXTRACT_CHARS:
        text = text[:MAX_EXTRACT_CHARS].rstrip() + "\n\n[... document truncated ...]"
    return text


def _decode_text(data: bytes) -> str:
    """Best-effort decode of a text-ish byte string."""
    for encoding in ("utf-8-sig", "utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(encoding)
        except (UnicodeDecodeError, LookupError):
            continue
    return data.decode("utf-8", errors="replace")


# ── Per-format extractors ───────────────────────────────────────────────────

def _extract_plaintext(data: bytes) -> str:
    return _decode_text(data)


def _extract_csv(data: bytes) -> str:
    text = _decode_text(data)
    rows = []
    reader = csv.reader(io.StringIO(text))
    for row in reader:
        cells = [c.strip() for c in row if c and c.strip()]
        if cells:
            rows.append(" | ".join(cells))
    return "\n".join(rows) if rows else text


def _extract_html(data: bytes) -> str:
    html = _decode_text(data)
    # Prefer trafilatura's main-content extraction; fall back to BeautifulSoup.
    try:
        import trafilatura
        extracted = trafilatura.extract(html)
        if extracted and extracted.strip():
            return extracted
    except Exception:
        pass
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(html, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        return soup.get_text(separator="\n")
    except Exception as exc:  # pragma: no cover - bs4 is a core dependency
        raise ExtractionError(f"Could not parse HTML: {exc}") from exc


def _extract_pdf(data: bytes) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise ExtractionError(
            "PDF support requires the 'pypdf' package. Install it with: pip install pypdf"
        ) from exc

    try:
        reader = PdfReader(io.BytesIO(data))
    except Exception as exc:
        raise ExtractionError(f"Could not read PDF: {exc}") from exc

    if getattr(reader, "is_encrypted", False):
        # Some PDFs are encrypted with an empty owner password; try that.
        try:
            reader.decrypt("")
        except Exception:
            raise ExtractionError(
                "This PDF is password-protected. Remove the password and try again."
            )

    parts = []
    for page in reader.pages:
        try:
            parts.append(page.extract_text() or "")
        except Exception:
            continue
    text = "\n\n".join(p for p in parts if p.strip())
    if not text.strip():
        raise ExtractionError(
            "No selectable text found in this PDF — it may be a scanned image. "
            "Try a PDF with real text, or paste the content into the From Text tab."
        )
    return text


def _extract_docx(data: bytes) -> str:
    try:
        import docx  # python-docx
    except ImportError as exc:
        raise ExtractionError(
            "Word support requires the 'python-docx' package. Install it with: pip install python-docx"
        ) from exc

    try:
        document = docx.Document(io.BytesIO(data))
    except Exception as exc:
        raise ExtractionError(f"Could not read Word document: {exc}") from exc

    parts = [p.text for p in document.paragraphs if p.text and p.text.strip()]
    # Include table cell text, which lives outside the paragraph stream.
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation  # python-pptx
    except ImportError as exc:
        raise ExtractionError(
            "PowerPoint support requires the 'python-pptx' package. Install it with: pip install python-pptx"
        ) from exc

    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as exc:
        raise ExtractionError(f"Could not read PowerPoint file: {exc}") from exc

    slides_text = []
    for idx, slide in enumerate(prs.slides, start=1):
        lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    run_text = "".join(run.text for run in para.runs)
                    if run_text.strip():
                        lines.append(run_text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
                    if cells:
                        lines.append(" | ".join(cells))
        # Speaker notes often carry the real narrative.
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f"(Notes: {notes})")
        if lines:
            slides_text.append(f"Slide {idx}:\n" + "\n".join(lines))
    return "\n\n".join(slides_text)


def _extract_xlsx(data: bytes) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise ExtractionError(
            "Excel support requires the 'openpyxl' package. Install it with: pip install openpyxl"
        ) from exc

    try:
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    except Exception as exc:
        raise ExtractionError(f"Could not read Excel file: {exc}") from exc

    sheets_text = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            sheets_text.append(f"Sheet: {ws.title}\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(sheets_text)


def _extract_rtf(data: bytes) -> str:
    """Best-effort RTF text extraction without an external dependency."""
    text = _decode_text(data)
    # Drop RTF groups/control words, keeping visible text.
    text = re.sub(r"\\'[0-9a-fA-F]{2}", "", text)          # hex-escaped chars
    text = re.sub(r"\\[a-zA-Z]+-?\d* ?", "", text)          # control words
    text = text.replace("{", "").replace("}", "")
    text = re.sub(r"\n{3,}", "\n\n", text)
    cleaned = text.strip()
    if not cleaned:
        raise ExtractionError("Could not extract text from this RTF file.")
    return cleaned


# ── Format registry ─────────────────────────────────────────────────────────
# Maps a lowercase extension (no dot) to its FormatSpec.

_FORMATS: dict[str, FormatSpec] = {
    "pdf":      FormatSpec("PDF", _extract_pdf, "pypdf"),
    "docx":     FormatSpec("Word (.docx)", _extract_docx, "docx"),
    "pptx":     FormatSpec("PowerPoint (.pptx)", _extract_pptx, "pptx"),
    "xlsx":     FormatSpec("Excel (.xlsx)", _extract_xlsx, "openpyxl"),
    "xlsm":     FormatSpec("Excel (.xlsm)", _extract_xlsx, "openpyxl"),
    "csv":      FormatSpec("CSV", _extract_csv, None),
    "tsv":      FormatSpec("TSV", _extract_csv, None),
    "txt":      FormatSpec("Text", _extract_plaintext, None),
    "md":       FormatSpec("Markdown", _extract_plaintext, None),
    "markdown": FormatSpec("Markdown", _extract_plaintext, None),
    "log":      FormatSpec("Log/Text", _extract_plaintext, None),
    "rtf":      FormatSpec("Rich Text (.rtf)", _extract_rtf, None),
    "html":     FormatSpec("HTML", _extract_html, None),
    "htm":      FormatSpec("HTML", _extract_html, None),
}


def _module_available(module_name: str | None) -> bool:
    if not module_name:
        return True
    import importlib.util
    return importlib.util.find_spec(module_name) is not None


def get_extension(filename: str) -> str:
    return filename.rsplit(".", 1)[1].lower() if "." in filename else ""


def is_supported(filename: str) -> bool:
    """True if we have a registered extractor for this file's extension."""
    return get_extension(filename) in _FORMATS


def supported_formats() -> list[dict]:
    """Report every known format and whether its dependency is installed.

    Returns a list of ``{"ext", "label", "available"}`` dicts (deduped by
    label), suitable for driving the compose UI's accepted-types list.
    """
    seen = set()
    out = []
    for ext, spec in _FORMATS.items():
        if spec.label in seen:
            continue
        seen.add(spec.label)
        out.append({
            "ext": ext,
            "label": spec.label,
            "available": _module_available(spec.requires),
        })
    return out


def accept_attribute() -> str:
    """Comma-separated ``.ext`` list for an <input accept="..."> attribute.

    Only includes formats whose dependency is currently importable.
    """
    exts = [
        f".{ext}"
        for ext, spec in _FORMATS.items()
        if _module_available(spec.requires)
    ]
    return ",".join(exts)


def extract_text(data: bytes, filename: str) -> str:
    """Extract plain text from ``data`` based on ``filename``'s extension.

    Raises :class:`ExtractionError` for unsupported types, missing optional
    dependencies, unreadable files, or documents with no extractable text.
    """
    if not filename:
        raise ExtractionError("No filename provided.")
    if len(data) > MAX_FILE_BYTES:
        raise ExtractionError("File is too large (max 16MB).")
    if not data:
        raise ExtractionError("The uploaded file is empty.")

    ext = get_extension(filename)
    spec = _FORMATS.get(ext)
    if spec is None:
        supported = ", ".join(sorted({s.label for s in _FORMATS.values()}))
        raise ExtractionError(
            f"Unsupported file type '.{ext or os.path.basename(filename)}'. "
            f"Supported formats: {supported}."
        )

    text = spec.extract(data)
    text = _truncate(text)
    if not text.strip():
        raise ExtractionError(
            f"No readable text was found in {filename}. "
            "The document may be empty, image-only, or corrupted."
        )
    return text
